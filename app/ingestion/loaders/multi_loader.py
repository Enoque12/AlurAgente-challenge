from pathlib import Path
import csv
import json

from bs4 import BeautifulSoup
from docx import Document as WordDocument
from openpyxl import load_workbook
from pptx import Presentation

from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
)


def load_file(file_path: str | Path) -> list[Document]:
    """
    Carrega um arquivo de acordo com sua extensão.
    """

    path = Path(file_path)
    extension = path.suffix.lower()

    if extension == ".pdf":
        documents = load_pdf(path)

    elif extension == ".docx":
        documents = load_docx(path)

    elif extension == ".pptx":
        documents = load_pptx(path)

    elif extension in {".md", ".txt"}:
        documents = load_text(path)

    elif extension in {".html", ".htm"}:
        documents = load_html(path)

    elif extension == ".csv":
        documents = load_csv(path)

    elif extension == ".json":
        documents = load_json(path)

    elif extension == ".xlsx":
        documents = load_xlsx(path)

    else:
        raise ValueError(
            f"Formato não suportado: {extension}"
        )

    # Normalizar metadados
    file_type = extension.lstrip(".")

    for document in documents:
        document.metadata["source"] = str(path)
        document.metadata["file_type"] = file_type

    return documents


def load_pdf(path: Path) -> list[Document]:
    """
    Carrega PDF utilizando PyPDF.
    """

    loader = PyPDFLoader(str(path))

    return loader.load()


def load_text(path: Path) -> list[Document]:
    """
    Carrega TXT ou Markdown.
    """

    loader = TextLoader(
        str(path),
        encoding="utf-8",
    )

    return loader.load()


def load_docx(path: Path) -> list[Document]:
    """
    Carrega texto de um documento Word.
    """

    document = WordDocument(str(path))

    paragraphs = []

    for paragraph in document.paragraphs:

        text = paragraph.text.strip()

        if text:
            paragraphs.append(text)

    content = "\n\n".join(paragraphs)

    if not content:
        return []

    return [
        Document(
            page_content=content,
            metadata={
                "source": str(path),
                "file_type": "docx",
            },
        )
    ]


def load_pptx(path: Path) -> list[Document]:
    """
    Extrai texto de todos os slides de um PowerPoint.
    """

    presentation = Presentation(str(path))

    documents: list[Document] = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        texts: list[str] = []

        for shape in slide.shapes:
            text = getattr(shape, "text", None)

            if not isinstance(text, str):
                continue

            text = text.strip()

            if text:
                texts.append(text)

        content = "\n".join(texts)

        if not content:
            continue

        documents.append(
            Document(
                page_content=content,
                metadata={
                    "source": str(path),
                    "file_type": "pptx",
                    "slide": slide_number,
                },
            )
        )

    return documents


def load_html(path: Path) -> list[Document]:
    """
    Extrai texto visível de uma página HTML.
    """

    html = path.read_text(
        encoding="utf-8",
    )

    soup = BeautifulSoup(
        html,
        "html.parser",
    )

    # Remover elementos que não representam
    # conteúdo útil para o RAG.
    for element in soup(
        ["script", "style", "noscript"]
    ):
        element.decompose()

    content = soup.get_text(
        separator="\n",
        strip=True,
    )

    if not content:
        return []

    return [
        Document(
            page_content=content,
            metadata={
                "source": str(path),
                "file_type": "html",
            },
        )
    ]


def load_csv(path: Path) -> list[Document]:
    """
    Carrega CSV.
    Cada linha é transformada em um Document.
    """

    documents = []

    with path.open(
        "r",
        encoding="utf-8-sig",
        newline="",
    ) as file:

        reader = csv.DictReader(file)

        for row_number, row in enumerate(
            reader,
            start=1,
        ):

            content = "\n".join(
                f"{key}: {value}"
                for key, value in row.items()
            )

            documents.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": str(path),
                        "file_type": "csv",
                        "row": row_number,
                    },
                )
            )

    return documents


def load_json(path: Path) -> list[Document]:
    """
    Carrega JSON.
    """

    with path.open(
        "r",
        encoding="utf-8",
    ) as file:

        data = json.load(file)

    content = json.dumps(
        data,
        ensure_ascii=False,
        indent=2,
    )

    return [
        Document(
            page_content=content,
            metadata={
                "source": str(path),
                "file_type": "json",
            },
        )
    ]


def load_xlsx(path: Path) -> list[Document]:
    """
    Carrega todas as planilhas de um arquivo Excel.
    """

    workbook = load_workbook(
        filename=str(path),
        data_only=True,
        read_only=True,
    )

    documents = []

    for worksheet in workbook.worksheets:

        rows = worksheet.iter_rows(
            values_only=True,
        )

        for row_number, row in enumerate(
            rows,
            start=1,
        ):

            values = [
                str(value)
                for value in row
                if value is not None
            ]

            if not values:
                continue

            content = " | ".join(values)

            documents.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": str(path),
                        "file_type": "xlsx",
                        "sheet": worksheet.title,
                        "row": row_number,
                    },
                )
            )

    workbook.close()

    return documents