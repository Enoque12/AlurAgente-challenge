from pathlib import Path
import json
from pptx import Presentation
from pptx.util import Inches

from docx import Document as WordDocument
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas


OUTPUT_DIR = Path("documents/test")


def create_markdown() -> None:
    path = OUTPUT_DIR / "teste.md"

    path.write_text(
        """# Documento de Teste

A NexusTech possui um benefício de transporte para colaboradores.

O valor definido atualmente é de 2.000 MZN por mês.
""",
        encoding="utf-8",
    )


def create_txt() -> None:
    path = OUTPUT_DIR / "teste.txt"

    path.write_text(
        """Documento de teste.

O horário de funcionamento da empresa é das 08:00 às 17:00.
""",
        encoding="utf-8",
    )


def create_csv() -> None:
    path = OUTPUT_DIR / "teste.csv"

    path.write_text(
        """produto,preco,stock
Laptop,45000,12
Monitor,18000,7
Teclado,2500,25
""",
        encoding="utf-8",
    )


def create_json() -> None:
    path = OUTPUT_DIR / "teste.json"

    data = {
        "empresa": "NexusTech",
        "departamento": "RH",
        "beneficio": "Seguro de saúde",
        "ativo": True,
    }

    path.write_text(
        json.dumps(
            data,
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )


def create_html() -> None:
    path = OUTPUT_DIR / "teste.html"

    path.write_text(
        """<!DOCTYPE html>
<html>
<head>
    <title>Documento de Teste</title>
</head>
<body>
    <h1>Política de Teste</h1>
    <p>A empresa oferece seguro de saúde aos colaboradores.</p>
</body>
</html>
""",
        encoding="utf-8",
    )


def create_docx() -> None:
    path = OUTPUT_DIR / "teste.docx"

    document = WordDocument()

    document.add_heading(
        "Documento de Teste",
        level=1,
    )

    document.add_paragraph(
        "A empresa oferece formação profissional "
        "aos seus colaboradores."
    )

    # python-docx aceita string como caminho.
    document.save(str(path))


def create_xlsx() -> None:
    path = OUTPUT_DIR / "teste.xlsx"

    workbook = Workbook()

    worksheet = workbook.active

    if worksheet is None:
        raise RuntimeError(
            "Não foi possível criar a planilha."
        )

    worksheet.title = "Benefícios"

    worksheet.append(
        ["Benefício", "Valor", "Periodicidade"]
    )

    worksheet.append(
        ["Transporte", "2000 MZN", "Mensal"]
    )

    worksheet.append(
        ["Alimentação", "1500 MZN", "Diário"]
    )

    workbook.save(str(path))


def create_pptx() -> None:
    path = OUTPUT_DIR / "teste.pptx"

    presentation = Presentation()

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    textbox = slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(1),
    )

    textbox.text_frame.text = "Documento de Teste"
    
    textbox2 = slide.shapes.add_textbox(
        Inches(1),
        Inches(2.5),
        Inches(8),
        Inches(2),
    )

    textbox2.text_frame.text = (
        "A NexusTech possui políticas internas "
        "de benefícios para colaboradores."
    )

    presentation.save(str(path))


def create_pdf() -> None:
    path = OUTPUT_DIR / "teste.pdf"

    pdf = canvas.Canvas(str(path))

    pdf.drawString(
        72,
        750,
        "Documento de Teste",
    )

    pdf.drawString(
        72,
        720,
        "A empresa oferece seguro de saúde.",
    )

    pdf.save()


def main() -> None:

    OUTPUT_DIR.mkdir(
        parents=True,
        exist_ok=True,
    )

    create_markdown()
    create_txt()
    create_csv()
    create_json()
    create_html()
    create_docx()
    create_xlsx()
    create_pptx()
    create_pdf()

    print(
        f"Documentos criados em: {OUTPUT_DIR}"
    )


if __name__ == "__main__":
    main()